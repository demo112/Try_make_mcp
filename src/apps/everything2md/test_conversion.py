import os
import sys
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from server import convert_to_markdown

def create_test_docx(path):
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test paragraph for DOCX conversion.')
    doc.save(path)

def create_test_xlsx(path):
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "Header 1"
    ws['B1'] = "Header 2"
    ws['A2'] = "Value 1"
    ws['B2'] = "Value 2"
    wb.save(path)

def create_test_pptx(path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    prs.save(path)

def run_tests():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    test_output_dir = os.path.join(base_dir, "test_outputs")
    
    if not os.path.exists(test_output_dir):
        os.makedirs(test_output_dir)

    # Define test files
    test_files = {
        "test.docx": create_test_docx,
        "test.xlsx": create_test_xlsx,
        "test.pptx": create_test_pptx
    }

    for filename, creator_func in test_files.items():
        source_path = os.path.join(test_output_dir, filename)
        output_path = os.path.join(test_output_dir, f"{filename}.md")
        
        print(f"Creating test file: {source_path}")
        creator_func(source_path)
        
        print(f"Converting {filename}...")
        result = convert_to_markdown(source_path, output_path)
        print(f"Result: {result}")
        
        if result == "Conversion successful" and os.path.exists(output_path):
            print(f"SUCCESS: {output_path} created.")
            # Optional: Read content to verify
            with open(output_path, 'r', encoding='utf-8') as f:
                print(f"Content preview:\n{f.read()[:200]}\n")
        else:
            print(f"FAILURE: Conversion failed for {filename}")

if __name__ == "__main__":
    run_tests()
